from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# File paths
input_path = r"C:\Users\music\OneDrive\Desktop\C2C_CRE_Website\Coast2CoastCREAcademyteaser.pptx"
output_path = r"C:\Users\music\OneDrive\Desktop\C2C_CRE_Website\Coast2CoastCREAcademyteaser_Enhanced.pptx"

# Load the original PowerPoint file
presentation = Presentation(input_path)

# Define color scheme and font
primary_color = RGBColor(0, 51, 102)  # Navy Blue
secondary_color = RGBColor(255, 255, 255)  # White
accent_color = RGBColor(255, 204, 0)  # Gold
font_name = "Calibri"

# Iterate through slides and redesign
for slide in presentation.slides:
    for shape in slide.shapes:
        if shape.has_text_frame:
            # Update text formatting
            for paragraph in shape.text_frame.paragraphs:
                paragraph.font.size = Pt(18)
                paragraph.font.name = font_name
                paragraph.font.color.rgb = primary_color
                paragraph.alignment = PP_ALIGN.LEFT

    # Customize the title slide
    if slide.shapes.title and "Title" in slide.shapes.title.text:
        slide.shapes.title.text = "Coast2Coast CRE Academy\nEmpowering the Future of Commercial Real Estate"
        slide.shapes.title.text_frame.paragraphs[0].font.size = Pt(32)
        slide.shapes.title.text_frame.paragraphs[0].font.bold = True
        slide.shapes.title.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

        # Set a bold background color
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = primary_color

# Save the redesigned presentation
presentation.save(output_path)

print(f"Enhanced PowerPoint saved as: {output_path}")
